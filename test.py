from pptx import Presentation
import pptx
import os, sys
from PIL import Image




def replace_image(prs, slide_index, old_image_path, new_image_path):
    slide = prs.slides[slide_index]
    print(slide.shapes)
    #print(slide.shape s[0].shape_type)
    for shape in slide.shapes:
        print(shape.shape_type)

        if shape.shape_type == 13:
            print('found image')
            slide.shapes._spTree.remove(shape._element)
             # Add a new image shape in the same position
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            slide.shapes.add_picture(new_image_path, left, top, width, height)
            break


 



def replace_image_in_presenation(prs, old_image_name, new_image_path):
    print('new image path', new_image_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == old_image_name and shape.shape_type == 13:
                print('found image')
                print(shape.name)
                print(slide.name)
                #save old shape attributes
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                print(width)
                print(height)
                new_pptx_img = pptx.parts.image.Image.from_file(new_image_path)
                slide_part, rId = shape.part, shape._element.blip_rId
                print(rId) 
                image_part = slide_part.related_part(rId)
                print(image_part)
               
                #print(new_image_path.blob)
                #new_pptx_img.dpi = shape.image.dpi
            
                image_part.blob = new_pptx_img._blob
                
                break
    prs.save('newTEST.pptx')
    

def parse_presentation(prs):
    print('START')
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                #print all shape attributes
                print('found image')
                print(shape.name)
                print(shape.image.size)
               
                
            


#presentation.save('modified_presentation.pptx')


#parse_presentation(presentation)
testPresentation = Presentation('test.pptx')
parse_presentation(testPresentation)






replace_image_in_presenation(testPresentation, 'Рисунок 29', 'p2.png')

# Save the modified presentation
